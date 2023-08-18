def ppt(valter1234567):

    import openai
    import os

# Set up OpenAI API credentials
    def generated_text():
        openai.api_key = "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"
    
        global query
        query =  valter1234567
# Set the prompt for the presentation
        prompt = '''Create a presentation on pitch deck on {} with slide number and its content in paragraph with 12 slides
                the each slide number should consist these 12 headings written below
                "Problem/ Opportunity"
            "Solution/ Technology"
            "Current Status/ Stage"
            "Product Demo"
            "Unique Value proposition"
            "Competitive Advantage"
            "Customer Segments & Market Size"
            "Channels (Paths proposed to be taken to reach your target customer segments for)"
            "Revenue streams"
            "Cost of the Project"
            "Means of Finance:"
            "Key Metrics & Validation"
            THANKYOU'''.format(query)

    # Generate the presentation text using GPT-3
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=prompt,
            max_tokens=1024,
            n=1,
            stop=None,
            temperature=0.6,
            #frequency_penalty = 0,
            presence_penalty = 1,
        )

        presentation_text = response.choices[0].text

        return presentation_text.lower()

    presentation_text = generated_text()

    def remove_all_occurrences(lst, value):
        while value in lst:
            lst.remove(value)
        return lst


    lines = presentation_text.split("\n")
    remove = ''
    lines = remove_all_occurrences(lines, remove)
    # Define empty lists to store the slide headings and contents
    slide_headings = []
    slide_contents = []

    # Loop over the lines and extract the slide headings and contents
    for i in range(len(lines)):
        line = lines[i].strip()
        if line.startswith("slide"):
            d = line.split(":")
            slide_headings.append(d[1])
            slide_contents.append(lines[i+1].strip())


    # Define an empty dictionary to store the slides
    slides = {}
    # Loop over the slide headings and contents and store them in the slides dictionary
    for i in range(len(slide_headings)):
        slides[slide_headings[i]] = slide_contents[i]


    import pptx
    from pptx.util import Inches

    #def display_template_options():
        #template_options = {
            #1: 'Template 1',
            #2: 'Template 2',
            #3: 'Template 3'
        #}
    
        #print("Available template options:")
        #for option, template_name in template_options.items():
            #print(f"{option}. {template_name}")
        
        #user_choice = input("Enter the number of the desired template: ")
        #return int(user_choice) if user_choice.isdigit() else None

    from pptx import Presentation
    # Display template options and get user choice
    #template_choice = None
    #while template_choice is None:
        #template_choice = display_template_options()

    # Use the template choice to generate the corresponding template
    template_path = "template_2.pptx"
    presentation = Presentation(template_path)

    # Create a new PowerPoint presentation
    prs = presentation

    # Loop over the slide headings and contents and add slides to the presentation
    for slide_heading, slide_content in slides.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        
        title.text = slide_heading.title()
        body.text = slide_content

    import requests
    from requests.structures import CaseInsensitiveDict

    import json

    QUERY_URL = "https://api.openai.com/v1/images/generations"

    def generate_image(prompt):
        model = "image-alpha-001"
        data = {
            "model": model,
            "prompt": prompt,
            "num_images": 1,
            "size": "512x512",
            "response_format": "url"
        }

        headers = CaseInsensitiveDict()
        headers["Content-Type"] = "application/json"
        api_key = "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"
        headers["Authorization"] = f"Bearer {api_key}"

        resp = requests.post(QUERY_URL, headers=headers, json=data)

        if resp.status_code != 200:
            raise ValueError("Failed to generate image")

        response_text = json.loads(resp.text)
        return response_text['data'][0]['url']


    def prompt():
        k= []
        li = [" Technology",
        "Stage",
        "Product Demo",
        "Unique Value proposition",
        "Competitive Advantage",
        "Customer Segments & Market Size",
        "Channels (Paths proposed to be taken to reach your target customer segments for)",
        "Revenue streams",
        "Cost of the Project",
        "Means of Finance:",
        "Key Metrics & Validation"
        ]

        
        for i in li:


            p = "image of {} for {}".format(i,query)
            image_url = generate_image(p)
            k.append(image_url)
            
        return k

    t = prompt()

    import urllib.request
    from pptx import Presentation
    from pptx.util import Inches
    import io

    # List of image URLs


    # Load the existing PowerPoint presentation
    presentation =prs

    # Loop through the image URLs
    for i, image_url in enumerate(t):
        # Download the image from the URL
        image_data = urllib.request.urlopen(image_url).read()

        # Get the slide at the corresponding index
        slide = presentation.slides[i]
        body = slide.placeholders[1]
    
        # Set the image width and height
        #image_width = Inches(4)
        #image_height = Inches(5)
        #slide_width = presentation.slide_width
        #image_left = slide_width - image_width -Inches(1)
        #image_top = Inches(2)
        image_width = Inches(4)
        image_height = Inches(5)
        body.left = Inches(1)
        body.top = Inches(2)
        body.width = Inches(7)
        body.height = Inches(4)
        # Calculate the position of the image to avoid overlap with the body placeholder
        image_left = body.left + body.width + Inches(0.5)  # Add a small offset to the right of the body placeholder
        image_top = body.top

        # Add the image to the slide
        # Add the image to the slide
        slide.shapes.add_picture(io.BytesIO(image_data), image_left, image_top, image_width, image_height)
        #slide.shapes.add_picture(image_data, image_left, image_top, image_width, image_height)
    
    # Save the modified presentation
    
    
    presentation.save("{}.pptx".format(query))

    return query




